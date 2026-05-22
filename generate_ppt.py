"""
거북목 판별 시스템 아키텍처 PPT 생성 스크립트
LLM(OpenAI / Gemini) 기반 거북목 여부 판별 시스템
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """PPT 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 슬라이드 1: 제목
    slide1_title_only = prs.slide_layouts[6]  # 빈 슬라이드
    slide1 = prs.slides.add_slide(slide1_title_only)
    
    # 배경색
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 100)
    
    # 제목
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = "거북목 판별 시스템 아키텍처"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 부제
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    
    p = subtitle_frame.paragraphs[0]
    p.text = "LLM(OpenAI / Gemini) 기반 거북목 여부 판별\n상품 20장 이미지 분석"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 슬라이드 2: 전체 데이터 흐름
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide2.shapes.title
    title.text = "전체 데이터 흐름"
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "사용자 (이미지 업로드)"
    
    points = [
        "프론트엔드 (Streamlit) - UI 및 이미지 처리",
        "백엔드 API (FastAPI) - 데이터 전송 및 저장",
        "특징 추출 (MediaPipe / OpenPose) - 포즈 및 자세 분석",
        "LLM 분석 (OpenAI / Gemini) - AI 기반 판별",
        "결과 반환 (거북목 여부 판별)"
    ]
    
    for i, point in enumerate(points[1:], 1):
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)
    
    # 슬라이드 3: 아키텍처 3계층
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide3.shapes.title
    title.text = "아키텍처 3계층"
    
    # 테이블 생성
    rows, cols = 4, 3
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    
    table_shape = slide3.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 헤더
    headers = ["계층", "역할", "기술"]
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(25, 45, 100)
        
        text_frame = cell.text_frame
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 데이터
    data = [
        ["1. 프론트엔드", "이미지 업로드 & UI", "Python, Streamlit"],
        ["2. 백엔드", "API 서버 & 데이터 처리", "FastAPI, SQLite"],
        ["3. AI/LLM", "지능형 분석 & 판별", "OpenAI(GPT-4o), Gemini 1.5"]
    ]
    
    colors = [RGBColor(41, 128, 185), RGBColor(39, 174, 96), RGBColor(155, 89, 182)]
    
    for row_idx, row_data in enumerate(data, 1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table_shape.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = colors[row_idx - 1]
            
            text_frame = cell.text_frame
            text_frame.paragraphs[0].font.size = Pt(12)
            text_frame.paragraphs[0].font.bold = True
            text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 슬라이드 4: 프론트엔드 (Streamlit)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide4.shapes.title
    title.text = "① 프론트엔드 (Streamlit)"
    
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        ("주요 기능:", [
            "이미지 업로드 (상품 20장)",
            "거북목 판별 웹 애플리케이션",
            "분석 결과 실시간 시각화"
        ]),
        ("분석 결과:", [
            "거북목 가능성: 78%",
            "평균 목 기울기: 26.4°",
            "기준 임계값: 20°",
            "분석 모델: Gemini-1.5-Flash",
            "분석 시간: 12.3s"
        ])
    ]
    
    for section_title, items in points:
        p = tf.add_paragraph()
        p.text = section_title
        p.level = 0
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(25, 45, 100)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(14)
    
    # 슬라이드 5: 백엔드 (FastAPI)
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide5.shapes.title
    title.text = "② 백엔드 (FastAPI + SQLite)"
    
    content = slide5.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        ("API 엔드포인트:", [
            "/upload - 이미지 업로드",
            "/analyze - 거북목 판별 요청",
            "/result/{id} - 분석 결과 조회"
        ]),
        ("전처리 & 특징 추출:", [
            "이미지 정규화 (Resizing, Normalization)",
            "자세/포즈 추출 (MediaPipe / OpenPose)",
            "특징 개선 (목 각도, 자세 지표)"
        ])
    ]
    
    for section_title, items in points:
        p = tf.add_paragraph()
        p.text = section_title
        p.level = 0
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(39, 174, 96)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(14)
    
    # 슬라이드 6: AI/LLM 계층
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide6.shapes.title
    title.text = "③ AI/LLM 계층 (인공지능)"
    
    content = slide6.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        ("LLM 모델:", [
            "OpenAI (GPT-4o)",
            "Google Gemini (1.5 Flash)"
        ]),
        ("프롬프트 엔지니어링:", [
            "각도 (angle) 분석",
            "어깨 (shoulder) 위치 분석",
            "정면/측면 비율 (ratio) 분석"
        ]),
        ("출력 결과:", [
            "거북목 여부 판정",
            "신뢰도 점수 (78%)",
            "개선 권장사항 및 지침"
        ])
    ]
    
    for section_title, items in points:
        p = tf.add_paragraph()
        p.text = section_title
        p.level = 0
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = RGBColor(155, 89, 182)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 1
            p.font.size = Pt(13)
    
    # 슬라이드 7: 기술 스택
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide7.shapes.title
    title.text = "기술 스택 & 배포"
    
    content = slide7.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    points = [
        ("Frontend:", "Python, Streamlit"),
        ("Backend:", "FastAPI, SQLite"),
        ("Pose Estimation:", "MediaPipe, OpenPose"),
        ("LLM:", "OpenAI API, Gemini API"),
        ("Deployment:", "Local / Cloud (선택 가능)"),
    ]
    
    for label, tech in points:
        p = tf.add_paragraph()
        p.text = f"{label} {tech}"
        p.level = 0
        p.font.size = Pt(15)
    
    # 슬라이드 8: 핵심 특징
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide8.shapes.title
    title.text = "핵심 특징"
    
    content = slide8.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        "✓ 실시간 거북목 판별 및 분석",
        "✓ LLM 기반 정확한 AI 판단",
        "✓ 사용자 친화적 웹 인터페이스",
        "✓ 빠른 처리 속도 (12.3s)",
        "✓ 높은 신뢰도 (78%)",
        "✓ 확장 가능한 모듈식 아키텍처",
        "✓ 로컬 및 클라우드 배포 지원"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(25, 45, 100)
    
    # 슬라이드 9: 감사합니다
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide9.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 45, 100)
    
    thank_you = slide9.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(3))
    tf = thank_you.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "감사합니다!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "거북목 판별 시스템 아키텍처"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 220, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # PPT 저장
    prs.save('presentation.pptx')
    print("✓ presentation.pptx 파일이 생성되었습니다!")

if __name__ == "__main__":
    create_presentation()
